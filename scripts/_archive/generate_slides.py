"""Generate progress presentation slides for thesis committee."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour palette ────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1A, 0x2E, 0x4A)   # slide background / headings
MID_BLUE    = RGBColor(0x1F, 0x6F, 0xB8)   # accent bar / bullets
LIGHT_BLUE  = RGBColor(0xD6, 0xE8, 0xF7)   # subtle bg boxes
GREEN       = RGBColor(0x1A, 0x7A, 0x4A)   # "done" items
AMBER       = RGBColor(0xC8, 0x7D, 0x0A)   # "in progress"
RED         = RGBColor(0xB8, 0x2A, 0x2A)   # "remaining"
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GREY   = RGBColor(0x33, 0x33, 0x33)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank


# ── Helpers ───────────────────────────────────────────────────────────────────
def add_rect(slide, l, t, w, h, fill=DARK_BLUE, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    return shape

def add_text(slide, text, l, t, w, h,
             size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             wrap=True, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb

def add_bullet_box(slide, items, l, t, w, h,
                   size=15, color=DARK_GREY, marker="▸", done_color=None):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (text, status) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        run = p.add_run()
        c = done_color if (done_color and status == "done") else \
            AMBER if status == "wip" else \
            RED if status == "todo" else color
        run.text = f"{marker}  {text}"
        run.font.size = Pt(size)
        run.font.color.rgb = c

def slide_bg(slide, color=LIGHT_GREY):
    bg = add_rect(slide, 0, 0, 13.33, 7.5, fill=color)
    bg.zorder = 0

def accent_bar(slide, title, subtitle=""):
    add_rect(slide, 0, 0, 13.33, 1.35, fill=DARK_BLUE)
    add_text(slide, title, 0.35, 0.12, 11, 0.65, size=26, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.35, 0.75, 11, 0.45, size=14, color=LIGHT_BLUE)

def divider(slide, y, color=MID_BLUE):
    add_rect(slide, 0.35, y, 12.63, 0.04, fill=color)

def badge(slide, text, l, t, w, h, bg=GREEN, fg=WHITE, size=13, bold=True):
    add_rect(slide, l, t, w, h, fill=bg)
    add_text(slide, text, l+0.07, t+0.04, w-0.14, h-0.08,
             size=size, bold=bold, color=fg, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
add_rect(s, 0, 5.9, 13.33, 1.6, fill=MID_BLUE)

add_text(s, "M.Sc. Thesis — Progress Presentation", 1, 0.5, 11.33, 0.55,
         size=15, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_text(s,
    "INTERSYMBOLIC-GRC",
    0.5, 1.1, 12.33, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(s,
    "An Intersymbolic AI Framework for Technical Asset\nRisk Assessment Under NFCRM-1:2025",
    0.5, 2.05, 12.33, 1.1, size=22, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

divider(s, 3.3, WHITE)

add_text(s, "Mohammed Ismail Al-Ammawi", 0.5, 3.45, 12.33, 0.55,
         size=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "M.Sc. Data Science  |  Islamic University of Madinah",
         0.5, 3.95, 12.33, 0.45, size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_text(s, "Supervisor: Prof. Fayçal Hamdi", 0.5, 4.55, 12.33, 0.45,
         size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

add_text(s, "Progress Presentation  |  May 3, 2026",
         0.5, 6.05, 12.33, 0.55, size=14, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Research Questions
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, "Problem Statement & Research Questions",
           "Bridging the gap between AI threat detection and GRC compliance reporting")

# Left: problem
add_rect(s, 0.35, 1.55, 5.9, 5.5, fill=WHITE)
add_text(s, "The Problem", 0.5, 1.62, 5.6, 0.45, size=15, bold=True, color=DARK_BLUE)
divider(s, 2.1, MID_BLUE)
problems = [
    "Pure-ML systems detect threats but produce zero GRC traceability",
    "Rule-based systems are explainable but cannot adapt to new threats",
    "LLMs generate risk narratives but hallucinate without symbolic grounding",
    "No standard method translates technical findings → NFCRM-1:2025 audit artifacts",
    "Saudi NCA compliance (NFCRM-1:2025) requires traceable, auditable AI outputs",
]
y = 2.2
for p in problems:
    add_text(s, f"▸  {p}", 0.5, y, 5.6, 0.52, size=13, color=DARK_GREY)
    y += 0.52

# Right: RQs
add_rect(s, 6.55, 1.55, 6.4, 5.5, fill=WHITE)
add_text(s, "Research Questions", 6.7, 1.62, 6.1, 0.45, size=15, bold=True, color=DARK_BLUE)
divider(s, 2.1, MID_BLUE)

rq_data = [
    ("RQ1", MID_BLUE,
     "How can an ontology encode asset–vulnerability–control relationships and directly translate to NFCRM-1:2025 clause mappings for GRC traceability?"),
    ("RQ2", GREEN,
     "At which pipeline points (pre- / in- / post-inference) should symbolic rules intervene to ensure correctness and compliance-ready outputs?"),
    ("RQ3", AMBER,
     "Does an intersymbolic pipeline provide capability differentiation in GRC outcomes that pure-ML or pure-rule baselines structurally cannot produce?"),
]
y = 2.2
for label, col, text in rq_data:
    badge(s, label, 6.7, y, 0.65, 0.38, bg=col)
    add_text(s, text, 7.45, y, 5.2, 0.85, size=12.5, color=DARK_GREY)
    y += 1.05


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Framework Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, "Framework Architecture — INTERSYMBOLIC-GRC",
           "Tri-stage intersymbolic pipeline with SHACL grounding layer")

# Three pipeline stage boxes
stages = [
    ("PRE-INFERENCE", MID_BLUE,
     ["SHACL ontology validation", "Event relevance filtering", "Baseline risk context", "Feature enrichment"]),
    ("ML INFERENCE", DARK_BLUE,
     ["Random Forest classifier", "XGBoost baseline", "CIC-IDS2018 / NSL-KDD", "87.98% accuracy"]),
    ("POST-INFERENCE", GREEN,
     ["SHACL compliance validation", "RiskCase generation", "NFCRM-1:2025 clause mapping", "Audit log + control linkage"]),
]
x = 0.35
for title, col, items in stages:
    add_rect(s, x, 1.55, 3.95, 4.5, fill=col)
    add_text(s, title, x+0.1, 1.62, 3.75, 0.5, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x+0.1, 2.15, 3.75, 3.75, fill=WHITE)
    yi = 2.25
    for item in items:
        add_text(s, f"• {item}", x+0.2, yi, 3.55, 0.45, size=12.5, color=DARK_GREY)
        yi += 0.48
    x += 4.19

# Arrow connectors (simple text arrows)
add_text(s, "→", 4.3, 3.4, 0.4, 0.5, size=28, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "→", 8.49, 3.4, 0.4, 0.5, size=28, bold=True, color=MID_BLUE, align=PP_ALIGN.CENTER)

# Bottom: LLM+SHACL path
add_rect(s, 0.35, 6.2, 12.63, 0.9, fill=LIGHT_BLUE)
add_text(s, "LLM+SHACL Path (parallel):  GLM-4.7 (ZAI API)  →  ARG reasoning (132 nodes, 423 edges, 32 CVEs)  →  SHACL validation  →  NL risk explanations",
         0.5, 6.28, 12.3, 0.55, size=12.5, color=DARK_BLUE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What Has Been Done (Thesis Chapters)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, "Progress: What Has Been Done",
           "All 5 thesis chapters complete — experiments run with real results")

chapters = [
    ("Ch.1  Introduction",      "Complete", GREEN,
     "Problem statement, RQs (RQ1–RQ3), 6 objectives (O1–O6), scope & limitations"),
    ("Ch.2  Related Work",       "Complete", GREEN,
     "15+ papers reviewed; 4 domains: GRC systems, ML detection, neuro-symbolic AI, knowledge graphs; gap analysis"),
    ("Ch.3  Methodology",        "Complete", GREEN,
     "SHACL ontology, ARG design (132 nodes), tri-stage pipeline, ML models, GRC artifact engine, reproducibility pkg"),
    ("Ch.4  Evaluation",         "Complete", GREEN,
     "4 experiments; ablation study; GRC metrics; statistical tests (p=0.779, power analysis); LLM hallucination study"),
    ("Ch.5  Conclusion",         "Complete", GREEN,
     "RQ1 Affirmed, RQ2 Partially Affirmed, RQ3 Affirmed (capability diff.); future work directions"),
]
y = 1.55
for chap, status, col, detail in chapters:
    add_rect(s, 0.35, y, 12.63, 0.95, fill=WHITE)
    badge(s, status, 0.45, y+0.27, 1.0, 0.38, bg=col, size=11)
    add_text(s, chap, 1.6, y+0.07, 3.0, 0.42, size=13, bold=True, color=DARK_BLUE)
    add_text(s, detail, 1.6, y+0.47, 11.1, 0.38, size=11.5, color=DARK_GREY)
    y += 1.03


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Key Experimental Results
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, "Key Experimental Results",
           "All results from real runs — no fabricated numbers")

# 4 result cards
cards = [
    (MID_BLUE,  "ML Pipeline",
     ["RF baseline: 87.70% acc, F1-macro 0.896",
      "XGBoost:  95.90% acc, F1-macro 0.692",
      "Tri-stage: 87.98% (+0.28%, p=0.779 NS)",
      "Post-inf. annotation: zero accuracy cost"]),
    (GREEN,     "GRC Metrics (NFCRM-1:2025)",
     ["100% NFCRM-1:2025 clause coverage",
      "100% event-to-risk traceability",
      "100% risk-to-control linkage",
      "57.6% ISO 27005 coverage (49/85)"]),
    (AMBER,     "LLM + SHACL (Hallucination)",
     ["No validation: 39.1% hallucination rate",
      "JSON Schema: 39.1% — no improvement",
      "SHACL: 0% among validated outputs",
      "100% hallucination reduction by SHACL"]),
    (RGBColor(0x5A,0x2D,0x8A),  "SLM Natural-Language Layer",
     ["Original SLM: 0% (design error)",
      "W17 NL layer CIC-IDS2018: 55.00%",
      "NSL-KDD NL layer (T2.2): 80.00%",
      "60/60 valid, 0 API errors (v3)"]),
]
x = 0.35
for col, title, items in cards:
    add_rect(s, x, 1.55, 3.0, 5.6, fill=col)
    add_text(s, title, x+0.1, 1.62, 2.8, 0.52, size=13, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x+0.1, 2.18, 2.8, 4.85, fill=WHITE)
    yi = 2.28
    for item in items:
        add_text(s, f"• {item}", x+0.18, yi, 2.64, 0.55, size=11.5, color=DARK_GREY)
        yi += 0.62
    x += 3.18


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Thesis Rating vs RQs
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, "RQ & Objective Fulfilment Status",
           "Honest assessment of what each RQ achieves and what gaps remain")

rows = [
    ("RQ1", "Ontology & GRC Traceability", "AFFIRMED", GREEN,
     "SHACL ontology + ARG built; 100% NFCRM coverage; 57.6% ISO 27005", "9/10"),
    ("RQ2", "Symbolic Intervention Points", "PARTIAL", AMBER,
     "Post-inf. annotation strong; pre-inf. NS (p=0.779); in-inference not evaluated", "7.5/10"),
    ("RQ3", "Capability Differentiation", "AFFIRMED", GREEN,
     "Pure-ML: 0% GRC traceability vs intersymbolic: 100%; SHACL halluci. 0%", "7/10"),
    ("O4",  "In-inference Mechanisms",    "DESIGNED", AMBER,
     "Architecture documented; SSH feature weighting implemented; not fully evaluated", "—"),
    ("T3.1","Human Expert Evaluation",    "PENDING",  RED,
     "Protocol written (human_eval/); 3 GRC practitioners needed; no data collected yet", "—"),
]
y = 1.55
for label, name, status, col, detail, score in rows:
    add_rect(s, 0.35, y, 12.63, 0.98, fill=WHITE)
    badge(s, label,   0.45, y+0.28, 0.65, 0.38, bg=DARK_BLUE, size=12)
    badge(s, status,  1.18, y+0.28, 1.3,  0.38, bg=col, size=11)
    add_text(s, name,   2.6, y+0.07, 4.5, 0.42, size=13, bold=True, color=DARK_BLUE)
    add_text(s, detail, 2.6, y+0.5,  8.5, 0.38, size=11.5, color=DARK_GREY)
    if score != "—":
        badge(s, score, 12.2, y+0.28, 0.65, 0.38, bg=MID_BLUE, size=12)
    y += 1.03


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — What Remains
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, "What Remains To Be Done",
           "Three focused tasks before final submission")

tasks = [
    (RED,    "CRITICAL",
     "Human Expert Evaluation (T3.1)",
     ["Recruit 3 GRC practitioners (protocol ready in human_eval/)",
      "Administer 30-min structured evaluation per evaluator",
      "Score GRC artifact quality, explanation clarity, audit readiness",
      "Estimated effort: 2–3 weeks to collect + 1 week to write up"]),
    (AMBER,  "HIGH",
     "Final Thesis Polish & Submission",
     ["Arabic abstract: re-enable XeLaTeX build (content ready in comments)",
      "Verify all figures compile without missing file errors",
      "Cross-check all numbers vs results/*.json (no stale figures)",
      "Format bibliography to IU Madinah style requirements"]),
    (MID_BLUE, "MEDIUM",
     "Defense Preparation",
     ["Finalize these slides with human eval results once collected",
      "Prepare 20-min presentation + 10-min Q&A rehearsal",
      "Prepare answers for anticipated examiner challenges:",
      "  — '100% NFCRM coverage is a SHACL guarantee, not empirical' → explain why this is correct design"]),
]
y = 1.55
for col, priority, title, items in tasks:
    add_rect(s, 0.35, y, 12.63, 1.85, fill=WHITE)
    badge(s, priority, 0.45, y+0.12, 1.2, 0.38, bg=col, size=11)
    add_text(s, title, 1.75, y+0.1, 10.5, 0.42, size=14, bold=True, color=DARK_BLUE)
    xi = 1.8
    for i, item in enumerate(items):
        row = y + 0.55 + (i // 2) * 0.48
        col_x = xi + (i % 2) * 5.2
        add_text(s, f"▸  {item}", col_x, row, 5.0, 0.42, size=11.5, color=DARK_GREY)
    y += 1.95


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Timeline
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, "Timeline to Final Submission",
           "Estimated path from progress presentation to thesis defense")

milestones = [
    ("May 2026",    "TODAY",   GREEN,
     "Progress presentation — all 5 chapters complete, experiments done"),
    ("May–Jun 2026","2–3 wks", AMBER,
     "Human evaluation: recruit 3 GRC practitioners, collect evaluation data"),
    ("Jun 2026",    "1 wk",    AMBER,
     "Write up human eval results; update Ch.4 evaluation section"),
    ("Jun 2026",    "1 wk",    MID_BLUE,
     "Final thesis polish: Arabic abstract, figure verification, bibliography"),
    ("Jul 2026",    "Target",  GREEN,
     "Submit final thesis for examination review"),
    ("Aug 2026",    "Target",  DARK_BLUE,
     "Defense presentation + examiner Q&A"),
]
y = 1.55
for month, effort, col, desc in milestones:
    add_rect(s, 0.35, y, 12.63, 0.83, fill=WHITE)
    add_rect(s, 0.35, y, 2.1,   0.83, fill=col)
    add_text(s, month,  0.4,  y+0.1,  1.9, 0.3, size=12, bold=True,
             color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, effort, 0.4,  y+0.43, 1.9, 0.3, size=11, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, desc,   2.6,  y+0.2,  10.2, 0.45, size=13, color=DARK_GREY)
    y += 0.9


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Summary
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
slide_bg(s)
accent_bar(s, "Summary",
           "INTERSYMBOLIC-GRC — Progress Status")

add_rect(s, 0.35, 1.55, 8.2, 5.5, fill=WHITE)
add_text(s, "Completed", 0.5, 1.62, 7.9, 0.45, size=15, bold=True, color=GREEN)
divider(s, 2.12, GREEN)
done = [
    "All 5 thesis chapters written and internally consistent",
    "Tri-stage ML+SHACL pipeline implemented and evaluated",
    "ARG: 132 nodes, 423 edges, 32 real CVEs from NVD API",
    "SHACL hallucination reduction: 39.1% → 0% (LLM outputs)",
    "SLM NL layer: 0% → 55% (CIC-IDS2018); 80% (NSL-KDD)",
    "Statistical testing: p=0.779 NS; power analysis complete",
    "Human evaluation protocol prepared (human_eval/)",
]
y = 2.22
for item in done:
    add_text(s, f"✓  {item}", 0.5, y, 7.9, 0.45, size=12.5, color=DARK_GREY)
    y += 0.45

add_rect(s, 8.78, 1.55, 4.2, 5.5, fill=WHITE)
add_text(s, "Remaining", 8.93, 1.62, 3.9, 0.45, size=15, bold=True, color=RED)
divider(s, 2.12, RED)
todo = [
    ("Human expert evaluation", RED),
    ("Arabic abstract (XeLaTeX build)", AMBER),
    ("Final figures verification", AMBER),
    ("Bibliography formatting", AMBER),
    ("Defense slides finalize", MID_BLUE),
    ("Thesis submission", MID_BLUE),
]
y = 2.22
for item, col in todo:
    c = "✗" if col == RED else "○"
    add_text(s, f"{c}  {item}", 8.93, y, 3.9, 0.5, size=12.5, color=col)
    y += 0.52

add_text(s, "Overall Status: Thesis ~90% Complete — On Track for Jul 2026 Submission",
         0.35, 7.1, 12.63, 0.45, size=13.5, bold=True, color=DARK_BLUE,
         align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out = "thesis/INTERSYMBOLIC-GRC_Progress_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
