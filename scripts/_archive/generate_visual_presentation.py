"""
Generate full visual presentation for INTERSYMBOLIC-GRC thesis committee.
Run: python scripts/generate_visual_presentation.py
"""
import json, os, textwrap
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.patches import FancyArrowPatch
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── paths ──────────────────────────────────────────────────────────────────
CHARTS = "thesis/figures/charts"
os.makedirs(CHARTS, exist_ok=True)
OUT    = "thesis/INTERSYMBOLIC-GRC_Visual_Presentation.pptx"

# ── colours ────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x1A, 0x2E, 0x4A)
MID_BLUE   = RGBColor(0x1F, 0x6F, 0xB8)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
GREEN      = RGBColor(0x1A, 0x7A, 0x4A)
AMBER      = RGBColor(0xC8, 0x7D, 0x0A)
RED        = RGBColor(0xB8, 0x2A, 0x2A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LGREY      = RGBColor(0xF2, 0xF4, 0xF7)
DGREY      = RGBColor(0x33, 0x33, 0x33)

# matplotlib palette
MC = {
    "blue":  "#1F6FB8",
    "dblue": "#1A2E4A",
    "green": "#1A7A4A",
    "amber": "#C87D0A",
    "red":   "#B82A2A",
    "lgrey": "#E8EDF2",
    "dgrey": "#333333",
    "purple":"#5A2D8A",
}

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.labelsize": 12,
    "axes.titlesize": 14,
    "axes.titleweight": "bold",
    "figure.facecolor": "white",
})

# ── load results ────────────────────────────────────────────────────────────
rf   = json.load(open("results/rf_baseline.json",              encoding="utf-8"))
xgb  = json.load(open("results/xgb_baseline.json",             encoding="utf-8"))
abl  = json.load(open("results/ablation_study_v2.json",        encoding="utf-8"))
grc  = json.load(open("results/grc_metrics.json",              encoding="utf-8"))
shacl= json.load(open("results/shacl_validation.json",         encoding="utf-8"))
slm  = json.load(open("results/slm_nl_classification.json",    encoding="utf-8"))
nsl  = json.load(open("results/slm_nslkdd_nl_classification.json", encoding="utf-8"))
expl = json.load(open("results/intersymbolic_explanations.json",   encoding="utf-8"))


# ════════════════════════════════════════════════════════════════════════════
# CHART GENERATORS
# ════════════════════════════════════════════════════════════════════════════

def save(fig, name):
    path = f"{CHARTS}/{name}.png"
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    print(f"  saved {path}")
    return path


# ── C1: ML performance comparison ──────────────────────────────────────────
def chart_ml_comparison():
    fig, axes = plt.subplots(1, 2, figsize=(12, 5))
    fig.suptitle("ML Pipeline Performance — CIC-IDS2018", fontsize=15, fontweight="bold",
                 color=MC["dblue"])

    labels   = ["RF Baseline", "XGBoost", "Tri-stage\n(Intersymbolic)"]
    accuracy = [87.70, 95.90, 87.98]
    f1_macro = [0.896, 0.692, 0.888]
    colors   = [MC["blue"], MC["purple"], MC["green"]]

    # Accuracy
    bars = axes[0].bar(labels, accuracy, color=colors, width=0.55, zorder=3)
    axes[0].set_ylim(80, 100)
    axes[0].set_ylabel("Accuracy (%)")
    axes[0].set_title("Accuracy")
    axes[0].yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    axes[0].set_facecolor(MC["lgrey"])
    for bar, val in zip(bars, accuracy):
        axes[0].text(bar.get_x() + bar.get_width()/2, val + 0.2,
                     f"{val:.2f}%", ha="center", va="bottom", fontsize=11, fontweight="bold")
    axes[0].annotate("+0.28% (p=0.779, NS)", xy=(2, 88.0), xytext=(1.4, 91),
                     arrowprops=dict(arrowstyle="->", color="grey"),
                     fontsize=9, color="grey", style="italic")

    # F1-Macro
    bars2 = axes[1].bar(labels, f1_macro, color=colors, width=0.55, zorder=3)
    axes[1].set_ylim(0, 1.1)
    axes[1].set_ylabel("F1-Macro")
    axes[1].set_title("F1-Macro (minority class quality)")
    axes[1].yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    axes[1].set_facecolor(MC["lgrey"])
    for bar, val in zip(bars2, f1_macro):
        axes[1].text(bar.get_x() + bar.get_width()/2, val + 0.01,
                     f"{val:.3f}", ha="center", va="bottom", fontsize=11, fontweight="bold")
    axes[1].annotate("RF chosen for pipeline:\nhigher F1-macro = better\nminority class recall",
                     xy=(0, 0.896), xytext=(0.6, 1.02),
                     arrowprops=dict(arrowstyle="->", color=MC["green"]),
                     fontsize=9, color=MC["green"])

    fig.tight_layout()
    return save(fig, "c1_ml_comparison")


# ── C2: Ablation study ──────────────────────────────────────────────────────
def chart_ablation():
    fig, ax = plt.subplots(figsize=(10, 5))
    fig.suptitle("Ablation Study — Pre/Post-Inference Symbolic Rules",
                 fontsize=14, fontweight="bold", color=MC["dblue"])

    configs  = ["ML Only\n(Baseline)", "Pre-Inference\nOnly", "Post-Inference\nOnly", "Full Tri-Stage\n(Intersymbolic)"]
    accuracy = [87.70, 87.98, 87.70, 87.98]
    grc_cov  = [0, 0, 100, 100]
    x = np.arange(len(configs))
    w = 0.35

    b1 = ax.bar(x - w/2, accuracy, w, label="Accuracy (%)", color=MC["blue"],   zorder=3)
    b2 = ax.bar(x + w/2, grc_cov,  w, label="NFCRM Coverage (%)", color=MC["green"], zorder=3)
    ax.set_ylim(0, 115)
    ax.set_ylabel("Score (%)")
    ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax.set_facecolor(MC["lgrey"])
    ax.set_xticks(x)
    ax.set_xticklabels(configs, fontsize=10)
    ax.legend(fontsize=11)

    for bar in b1:
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.5,
                f"{bar.get_height():.2f}%", ha="center", fontsize=9, fontweight="bold")
    for bar in b2:
        v = bar.get_height()
        col = MC["green"] if v == 100 else MC["red"]
        ax.text(bar.get_x()+bar.get_width()/2, v+0.5,
                f"{v:.0f}%", ha="center", fontsize=9, fontweight="bold", color=col)

    ax.annotate("Post-inference adds\n100% GRC coverage\nat zero accuracy cost",
                xy=(2.18, 100), xytext=(2.5, 108),
                arrowprops=dict(arrowstyle="->", color=MC["green"]),
                fontsize=9, color=MC["green"], fontweight="bold")
    fig.tight_layout()
    return save(fig, "c2_ablation")


# ── C3: GRC metrics donut ───────────────────────────────────────────────────
def chart_grc_metrics():
    fig, axes = plt.subplots(1, 3, figsize=(13, 5))
    fig.suptitle("GRC Metrics — NFCRM-1:2025 Compliance",
                 fontsize=14, fontweight="bold", color=MC["dblue"])

    metrics = [
        ("NFCRM-1:2025\nClause Coverage", 100, 0),
        ("Event-to-Risk\nTraceability",    100, 0),
        ("ISO 27005\nCoverage",            57.6, 42.4),
    ]
    for ax, (title, val, rest) in zip(axes, metrics):
        col = MC["green"] if val == 100 else MC["amber"]
        wedges, _ = ax.pie(
            [val, rest] if rest > 0 else [val],
            colors=[col, "#E8EDF2"] if rest > 0 else [col],
            startangle=90, counterclock=False,
            wedgeprops=dict(width=0.45, edgecolor="white", linewidth=2)
        )
        ax.text(0, 0, f"{val:.0f}%", ha="center", va="center",
                fontsize=22, fontweight="bold", color=col)
        ax.set_title(title, fontsize=12, fontweight="bold", color=MC["dblue"], pad=10)

    axes[2].text(0, -0.72, "49/85 RiskCases\nInfiltration §6.5 pending",
                 ha="center", fontsize=9, color=MC["amber"], style="italic")
    axes[0].text(0, -0.72, "SHACL-enforced\nguarantee",
                 ha="center", fontsize=9, color=MC["green"], style="italic")
    axes[1].text(0, -0.72, "Full audit trail\nEvent→RiskCase→Control",
                 ha="center", fontsize=9, color=MC["green"], style="italic")
    fig.tight_layout()
    return save(fig, "c3_grc_metrics")


# ── C4: SHACL hallucination reduction ──────────────────────────────────────
def chart_shacl():
    fig, ax = plt.subplots(figsize=(10, 5))
    fig.suptitle("SHACL Hallucination Reduction — LLM+SHACL Experiment",
                 fontsize=14, fontweight="bold", color=MC["dblue"])

    levels  = ["No Validation", "JSON Schema Only", "Full SHACL"]
    rates   = [39.1, 39.1, 0.0]
    cols    = [MC["red"], MC["amber"], MC["green"]]
    bars    = ax.bar(levels, rates, color=cols, width=0.5, zorder=3)
    ax.set_ylim(0, 55)
    ax.set_ylabel("Hallucination Rate (%)")
    ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax.set_facecolor(MC["lgrey"])

    for bar, val, col in zip(bars, rates, cols):
        label = f"{val:.1f}%" if val > 0 else "0%  (100% reduction)"
        ax.text(bar.get_x()+bar.get_width()/2,
                val + 0.8 if val > 0 else 0.8,
                label, ha="center", fontsize=12,
                fontweight="bold", color=col)

    ax.annotate("JSON Schema:\nzero improvement\nover no validation",
                xy=(1, 39.1), xytext=(1.25, 48),
                arrowprops=dict(arrowstyle="->", color=MC["amber"]),
                fontsize=9, color=MC["amber"])
    ax.annotate("SHACL: semantic\ngrounding catches\nwhat syntax cannot",
                xy=(2, 0), xytext=(1.55, 20),
                arrowprops=dict(arrowstyle="->", color=MC["green"]),
                fontsize=9, color=MC["green"], fontweight="bold")

    ax.text(0.5, -0.12,
            "N=64 parseable outputs  |  25 hallucinations found and blocked by SHACL",
            ha="center", transform=ax.transAxes,
            fontsize=9, color="grey", style="italic")
    fig.tight_layout()
    return save(fig, "c4_shacl")


# ── C5: SLM NL conversion — CIC-IDS2018 ───────────────────────────────────
def chart_slm_cicids():
    fig, axes = plt.subplots(1, 2, figsize=(13, 5))
    fig.suptitle("SLM Natural-Language Classification — CIC-IDS2018",
                 fontsize=14, fontweight="bold", color=MC["dblue"])

    # Left: before/after
    ax = axes[0]
    categories = ["Without NL\nConversion\n(design error)", "With NL\nConversion\n(W17 v3)"]
    values     = [0, 55]
    cols       = [MC["red"], MC["green"]]
    bars       = ax.bar(categories, values, color=cols, width=0.45, zorder=3)
    ax.set_ylim(0, 80)
    ax.set_ylabel("Accuracy (%)")
    ax.set_title("Before vs After NL Conversion")
    ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax.set_facecolor(MC["lgrey"])
    for bar, val, col in zip(bars, values, cols):
        ax.text(bar.get_x()+bar.get_width()/2, val+1,
                f"{val}%", ha="center", fontsize=14, fontweight="bold", color=col)
    ax.annotate("+55pp\nimprovement", xy=(0.5, 27), fontsize=12,
                ha="center", color=MC["green"], fontweight="bold")
    ax.text(0.5, -0.15, "60/60 valid responses  |  0 API errors",
            ha="center", transform=ax.transAxes, fontsize=9, color="grey", style="italic")

    # Right: per-class breakdown
    ax2   = axes[1]
    cls   = list(slm["per_class_accuracy"].keys())
    vals  = [v*100 for v in slm["per_class_accuracy"].values()]
    cols2 = [MC["green"] if v >= 60 else MC["amber"] for v in vals]
    bars2 = ax2.barh(cls, vals, color=cols2, zorder=3)
    ax2.set_xlim(0, 100)
    ax2.set_xlabel("Accuracy (%)")
    ax2.set_title("Per-Class Accuracy")
    ax2.xaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax2.set_facecolor(MC["lgrey"])
    for bar, val in zip(bars2, vals):
        ax2.text(val+1, bar.get_y()+bar.get_height()/2,
                 f"{val:.0f}%", va="center", fontsize=11, fontweight="bold")
    ax2.axvline(55, color=MC["blue"], linestyle="--", linewidth=1.5, label="Overall avg 55%")
    ax2.legend(fontsize=9)

    fig.tight_layout()
    return save(fig, "c5_slm_cicids")


# ── C6: SLM NSL-KDD ─────────────────────────────────────────────────────────
def chart_slm_nslkdd():
    fig, axes = plt.subplots(1, 2, figsize=(13, 5))
    fig.suptitle("SLM Natural-Language Classification — NSL-KDD (T2.2)",
                 fontsize=14, fontweight="bold", color=MC["dblue"])

    # Left: comparison bar
    ax = axes[0]
    datasets = ["CIC-IDS2018\n(W17 v3)", "NSL-KDD\n(T2.2)"]
    accs     = [55, 74]
    cols     = [MC["blue"], MC["green"]]
    bars     = ax.bar(datasets, accs, color=cols, width=0.45, zorder=3)
    ax.set_ylim(0, 100)
    ax.set_ylabel("Accuracy (%)")
    ax.set_title("SLM Accuracy Across Datasets")
    ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax.set_facecolor(MC["lgrey"])
    for bar, val in zip(bars, accs):
        ax.text(bar.get_x()+bar.get_width()/2, val+1,
                f"{val}%", ha="center", fontsize=14, fontweight="bold")
    ax.axhline(55, color="grey", linestyle=":", linewidth=1)

    # Right: NSL-KDD per class
    ax2  = axes[1]
    cls  = list(nsl["per_class_accuracy"].keys())
    vals = [v*100 for v in nsl["per_class_accuracy"].values()]
    cols2= [MC["green"] if v >= 80 else MC["amber"] if v >= 50 else MC["red"] for v in vals]
    bars2= ax2.barh(cls, vals, color=cols2, zorder=3)
    ax2.set_xlim(0, 115)
    ax2.set_xlabel("Accuracy (%)")
    ax2.set_title("NSL-KDD Per-Class Accuracy")
    ax2.xaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax2.set_facecolor(MC["lgrey"])
    for bar, val in zip(bars2, vals):
        tag = " (perfect)" if val == 100 else ""
        ax2.text(val+1, bar.get_y()+bar.get_height()/2,
                 f"{val:.0f}%{tag}", va="center", fontsize=11, fontweight="bold")
    ax2.axvline(74, color=MC["blue"], linestyle="--", linewidth=1.5, label="Overall avg 74%")
    ax2.legend(fontsize=9)
    ax2.text(0.5, -0.15, "50/50 valid  |  0 API errors",
             ha="center", transform=ax2.transAxes, fontsize=9, color="grey", style="italic")

    fig.tight_layout()
    return save(fig, "c6_slm_nslkdd")


# ── C7: ARG stats ────────────────────────────────────────────────────────────
def chart_arg():
    fig, axes = plt.subplots(1, 2, figsize=(13, 5))
    fig.suptitle("Asset Relationship Graph (ARG) — Multi-Source Knowledge Graph",
                 fontsize=14, fontweight="bold", color=MC["dblue"])

    # Left: node breakdown
    ax = axes[0]
    node_types = ["Assets\n(CMDB)", "Vulnerabilities\n(CVE/NVD)", "ATT&CK\nTechniques",
                  "NFCRM\nControls", "Network\nEvents", "Other"]
    counts     = [50, 32, 20, 15, 10, 5]
    cols       = [MC["blue"], MC["red"], MC["amber"], MC["green"], MC["purple"], "grey"]
    wedges, texts, autotexts = ax.pie(
        counts, labels=node_types, colors=cols,
        autopct="%1.0f%%", startangle=140,
        pctdistance=0.75,
        wedgeprops=dict(edgecolor="white", linewidth=2)
    )
    for t in autotexts:
        t.set_fontsize(9)
        t.set_fontweight("bold")
    ax.set_title(f"132 Total Nodes — 5 Data Sources", fontsize=12, fontweight="bold",
                 color=MC["dblue"])

    # Right: source stats
    ax2 = axes[1]
    sources = ["CMDB\n(Assets)", "NVD API\n(CVEs)", "MITRE\nATT&CK", "NFCRM\nControls", "Network\nFlows"]
    vals    = [50, 32, 20, 15, 15]
    edge_ct = [150, 120, 80, 50, 23]
    x = np.arange(len(sources))
    w = 0.38
    b1 = ax2.bar(x-w/2, vals,    w, label="Nodes", color=MC["blue"],  zorder=3)
    b2 = ax2.bar(x+w/2, edge_ct, w, label="Edges", color=MC["amber"], zorder=3)
    ax2.set_ylabel("Count")
    ax2.set_title("Nodes & Edges per Source", fontsize=12, fontweight="bold", color=MC["dblue"])
    ax2.set_xticks(x)
    ax2.set_xticklabels(sources, fontsize=9)
    ax2.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax2.set_facecolor(MC["lgrey"])
    ax2.legend(fontsize=10)
    for bar in b1:
        ax2.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1,
                 str(int(bar.get_height())), ha="center", fontsize=9)
    for bar in b2:
        ax2.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1,
                 str(int(bar.get_height())), ha="center", fontsize=9)
    ax2.text(0.5, -0.15, "Total: 132 nodes  |  423 edges  |  32 real CVEs from NVD",
             ha="center", transform=ax2.transAxes, fontsize=9, color=MC["dblue"],
             fontweight="bold")

    fig.tight_layout()
    return save(fig, "c7_arg")


# ── C8: Effort timeline (commit activity) ───────────────────────────────────
def chart_effort():
    weeks  = ["W06\nFeb 9", "W07\nFeb 16", "W08\nFeb 23", "W09\nMar 1",
              "W10\nMar 8", "W11\nMar 15", "W12\nMar 22", "W13\nMar 29",
              "W14\nApr 5", "W15\nApr 12", "W17\nApr 26", "W18\nMay 3"]
    commits= [19, 31, 2, 48, 31, 137, 128, 54, 108, 5, 8, 2]
    phases = ["Foundation", "Foundation", "Foundation", "Data+Exp",
              "Evaluation", "PEAK Eval", "PEAK Eval", "Stabilise",
              "Quality", "Final", "Final", "Present"]
    phase_cols = {
        "Foundation": MC["blue"], "Data+Exp": MC["purple"],
        "Evaluation": MC["amber"], "PEAK Eval": MC["red"],
        "Stabilise": MC["amber"], "Quality": MC["green"],
        "Final": MC["green"], "Present": MC["dblue"],
    }
    cols = [phase_cols[p] for p in phases]

    fig, ax = plt.subplots(figsize=(14, 6))
    fig.suptitle("Weekly Commit Activity — 573 Total Commits Over 12 Weeks",
                 fontsize=14, fontweight="bold", color=MC["dblue"])

    bars = ax.bar(weeks, commits, color=cols, zorder=3, width=0.7)
    ax.set_ylabel("Commits")
    ax.yaxis.grid(True, linestyle="--", alpha=0.4, zorder=0)
    ax.set_facecolor(MC["lgrey"])

    for bar, val in zip(bars, commits):
        if val > 5:
            ax.text(bar.get_x()+bar.get_width()/2, val+2,
                    str(val), ha="center", fontsize=9, fontweight="bold")

    # phase legend
    unique_phases = list(dict.fromkeys(phases))
    handles = [mpatches.Patch(color=phase_cols[p], label=p) for p in unique_phases]
    ax.legend(handles=handles, loc="upper right", fontsize=9, ncol=2)

    # annotations
    ax.annotate("Peak: 137\nEval chapter\n+ verification",
                xy=(5, 137), xytext=(5.5, 120),
                arrowprops=dict(arrowstyle="->", color=MC["red"]),
                fontsize=8, color=MC["red"])
    ax.annotate("W1-W22\nquality fixes",
                xy=(8, 108), xytext=(8.5, 118),
                arrowprops=dict(arrowstyle="->", color=MC["green"]),
                fontsize=8, color=MC["green"])

    ax.text(0.5, -0.15,
            "52 active days  |  195+ PRs merged  |  Feb 9 – May 3, 2026",
            ha="center", transform=ax.transAxes, fontsize=10,
            color=MC["dblue"], fontweight="bold")
    fig.tight_layout()
    return save(fig, "c8_effort")


# ── C9: Capability matrix ───────────────────────────────────────────────────
def chart_capability():
    fig, ax = plt.subplots(figsize=(12, 5))
    fig.suptitle("Capability Matrix — What Each Approach Produces",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    ax.axis("off")

    approaches = ["Pure-Rule\nSystem", "Pure-ML\nSystem", "ML+SHACL\n(This work)", "LLM+SHACL\n(This work)"]
    capabilities = [
        "Adapts to\nnew threats",
        "Standards\nclause mapping",
        "Audit-ready\nRiskCase",
        "NL risk\nexplanation",
        "Hallucination\nprevention",
        "NFCRM-1:2025\ncoverage",
    ]
    matrix = [
        # Rule  ML   ML+S  LLM+S
        [False, True,  True,  True ],   # adapts
        [True,  False, True,  True ],   # standards
        [False, False, True,  True ],   # riskcase
        [False, False, False, True ],   # NL explanation
        [False, False, True,  True ],   # hallucination
        [False, False, True,  True ],   # NFCRM
    ]

    col_w = 0.18
    row_h = 0.13
    x_start = 0.25
    y_start = 0.85

    # headers
    for j, app in enumerate(approaches):
        ax.text(x_start + j*col_w + col_w/2, y_start + 0.04,
                app, ha="center", va="bottom", fontsize=10,
                fontweight="bold", color=MC["dblue"],
                transform=ax.transAxes)

    # row labels + cells
    for i, cap in enumerate(capabilities):
        y = y_start - i*row_h
        ax.text(0.22, y - row_h/2, cap, ha="right", va="center",
                fontsize=10, transform=ax.transAxes, color=MC["dgrey"])
        for j, has in enumerate(matrix[i]):
            x = x_start + j*col_w + col_w/2
            bg  = MC["green"] if has else "#FFDDDD"
            sym = "YES" if has else "NO"
            fc  = "white" if has else MC["red"]
            rect = plt.Rectangle((x - col_w/2.2, y - row_h*0.85),
                                  col_w*0.9, row_h*0.78,
                                  transform=ax.transAxes,
                                  facecolor=bg, edgecolor="white",
                                  linewidth=2, clip_on=False)
            ax.add_patch(rect)
            ax.text(x, y - row_h*0.45, sym, ha="center", va="center",
                    fontsize=9, fontweight="bold", color=fc,
                    transform=ax.transAxes)

    fig.tight_layout()
    return save(fig, "c9_capability")


# ── C10: 4-layer pipeline visual ────────────────────────────────────────────
def chart_pipeline():
    fig, ax = plt.subplots(figsize=(14, 5))
    fig.suptitle("4-Layer Intersymbolic Pipeline — End-to-End Flow",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    ax.axis("off")

    layers = [
        ("Layer 1\nRaw Features",   MC["blue"],   "CICFlowMeter\nnumeric columns"),
        ("Layer 2\nSymbolic Rules", MC["amber"],  "Pre-inference\nSHACL validation\n& filtering"),
        ("Layer 3\nML + SHACL",     MC["green"],  "Random Forest\n→ RiskCase\n→ NFCRM clause"),
        ("Layer 4\nNL Explanation", MC["purple"], "GLM-4.7\n(ZAI API)\nGRC narrative"),
    ]
    extras = [
        "16.2M flows\n(CIC-IDS2018)",
        "+0.28% accuracy\n(NS, p=0.779)",
        "100% NFCRM\ncoverage",
        "All 5 flows\n(0 errors)",
    ]

    n = len(layers)
    for i, ((title, col, desc), extra) in enumerate(zip(layers, extras)):
        x = 0.07 + i * 0.235
        # box
        rect = plt.Rectangle((x, 0.15), 0.20, 0.70,
                              transform=ax.transAxes, facecolor=col,
                              edgecolor="white", linewidth=3, clip_on=False)
        ax.add_patch(rect)
        ax.text(x+0.10, 0.78, title, ha="center", va="center",
                fontsize=11, fontweight="bold", color="white",
                transform=ax.transAxes)
        ax.text(x+0.10, 0.48, desc, ha="center", va="center",
                fontsize=9, color="white", linespacing=1.5,
                transform=ax.transAxes)
        # extra label below
        ax.text(x+0.10, 0.08, extra, ha="center", va="center",
                fontsize=8.5, color=col, fontweight="bold", style="italic",
                transform=ax.transAxes)
        # arrow
        if i < n-1:
            ax.annotate("", xy=(x+0.235, 0.50), xytext=(x+0.205, 0.50),
                        xycoords="axes fraction", textcoords="axes fraction",
                        arrowprops=dict(arrowstyle="-|>", color="grey",
                                        lw=2.5, mutation_scale=20))

    fig.tight_layout()
    return save(fig, "c10_pipeline")


# ════════════════════════════════════════════════════════════════════════════
# GENERATE ALL CHARTS
# ════════════════════════════════════════════════════════════════════════════
print("Generating charts...")
charts = {
    "pipeline":    chart_pipeline(),
    "ml":          chart_ml_comparison(),
    "ablation":    chart_ablation(),
    "grc":         chart_grc_metrics(),
    "shacl":       chart_shacl(),
    "slm_cicids":  chart_slm_cicids(),
    "slm_nslkdd":  chart_slm_nslkdd(),
    "arg":         chart_arg(),
    "effort":      chart_effort(),
    "capability":  chart_capability(),
}
print(f"  {len(charts)} charts ready.\n")


# ════════════════════════════════════════════════════════════════════════════
# BUILD PPTX
# ════════════════════════════════════════════════════════════════════════════
BLANK = None
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

def bg(slide, col=LGREY):
    s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = col
    s.line.fill.background()

def header(slide, title, subtitle=""):
    s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.25))
    s.fill.solid(); s.fill.fore_color.rgb = DARK_BLUE
    s.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.08), Inches(12.5), Inches(0.6))
    tf = tb.text_frame; p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = title
    r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = WHITE
    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.3), Inches(0.72), Inches(12.5), Inches(0.4))
        tf2 = tb2.text_frame; p2 = tf2.paragraphs[0]
        r2 = p2.add_run(); r2.text = subtitle
        r2.font.size = Pt(12); r2.font.color.rgb = LIGHT_BLUE

def add_img(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def add_txt(slide, text, l, t, w, h, size=13, bold=False,
            color=DGREY, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

def badge(slide, text, l, t, w, h, bg_col=GREEN, fg=WHITE, size=12):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = bg_col; s.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(l+0.05), Inches(t+0.03), Inches(w-0.1), Inches(h-0.06))
    tf = tb.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = True; r.font.color.rgb = fg

print("Building presentation slides...")

# ── S1: Title ─────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
sh = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
sh.fill.solid(); sh.fill.fore_color.rgb = DARK_BLUE; sh.line.fill.background()
add_img(s, charts["pipeline"], 0.2, 3.5, 12.9, 3.6)

# dark overlay on image area
ov = s.shapes.add_shape(1, 0, Inches(3.4), prs.slide_width, Inches(3.7))
ov.fill.solid(); ov.fill.fore_color.rgb = DARK_BLUE
ov.line.fill.background()
from pptx.util import Pt as _Pt
from pptx.dml.color import RGBColor as _RGB
from pptx.oxml.ns import qn
# set transparency via XML
sp = ov.fill._xPr
sp.find(qn("a:solidFill")).find(qn("a:srgbClr")).set("val", "1A2E4A")

add_img(s, charts["pipeline"], 0.15, 3.55, 13.0, 3.6)

add_txt(s, "M.Sc. Progress Presentation  |  May 3, 2026",
        0.5, 0.25, 12.33, 0.5, size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_txt(s, "INTERSYMBOLIC-GRC",
        0.5, 0.75, 12.33, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_txt(s, "An Intersymbolic AI Framework for Technical Asset Risk Assessment",
        0.5, 1.70, 12.33, 0.55, size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_txt(s, "Under NFCRM-1:2025",
        0.5, 2.20, 12.33, 0.5, size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_txt(s, "Mohammed Ismail Al-Ammawi  |  Islamic University of Madinah",
        0.5, 2.85, 12.33, 0.5, size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_txt(s, "Supervisor: Prof. Faycal Hamdi",
        0.5, 3.2, 12.33, 0.4, size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
print("  Slide 1 done")

# ── S2: Problem + RQs ─────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Problem Statement & Research Questions",
              "Why current GRC systems fail and what this work addresses")

# Problem box
pb = s.shapes.add_shape(1, Inches(0.3), Inches(1.4), Inches(5.8), Inches(5.8))
pb.fill.solid(); pb.fill.fore_color.rgb = WHITE; pb.line.fill.background()
add_txt(s, "The Gap", 0.4, 1.45, 5.6, 0.4, size=14, bold=True, color=DARK_BLUE)

problems = [
    ("Pure-ML systems", "detect threats but produce ZERO GRC traceability"),
    ("Rule-based systems", "are explainable but cannot adapt to new threats"),
    ("LLMs", "generate risk text but hallucinate without grounding"),
    ("No existing system", "translates findings -> NFCRM-1:2025 audit reports"),
    ("Saudi NCA mandate", "requires traceable, auditable AI outputs"),
]
y = 1.95
for bold_part, rest in problems:
    tb = s.shapes.add_textbox(Inches(0.45), Inches(y), Inches(5.5), Inches(0.52))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r1 = p.add_run(); r1.text = bold_part + " — "
    r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = MID_BLUE
    r2 = p.add_run(); r2.text = rest
    r2.font.size = Pt(12); r2.font.color.rgb = DGREY
    y += 0.55

# RQ boxes
rq_data = [
    (RGBColor(0x1F,0x6F,0xB8), "RQ1",
     "How can an ontology encode asset-vulnerability-control relationships and map directly to NFCRM-1:2025 clauses?"),
    (RGBColor(0x1A,0x7A,0x4A), "RQ2",
     "At which pipeline points should symbolic rules intervene to ensure correctness and compliance-ready outputs?"),
    (RGBColor(0xC8,0x7D,0x0A), "RQ3",
     "Does the intersymbolic pipeline provide GRC capability differentiation that pure-ML/pure-rule cannot produce?"),
]
y = 1.45
for col, label_txt, q_txt in rq_data:
    rb = s.shapes.add_shape(1, Inches(6.35), Inches(y), Inches(6.65), Inches(1.65))
    rb.fill.solid(); rb.fill.fore_color.rgb = WHITE; rb.line.fill.background()
    badge(s, label_txt, 6.4, y+0.12, 0.75, 0.42, bg_col=col, size=13)
    add_txt(s, q_txt, 7.25, y+0.08, 5.65, 1.3, size=12, color=DGREY)
    y += 1.75
print("  Slide 2 done")

# ── S3: Pipeline architecture ─────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Framework Architecture — 4-Layer Intersymbolic Pipeline",
              "ML+SHACL path + LLM+SHACL path working in parallel")
add_img(s, charts["pipeline"], 0.2, 1.35, 12.9, 5.2)
print("  Slide 3 done")

# ── S4: ML results ────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Experiment 1 — ML Pipeline Performance",
              "CIC-IDS2018 (2,122 test samples) | RF chosen for tri-stage (superior F1-macro)")
add_img(s, charts["ml"], 0.2, 1.35, 12.9, 5.5)
print("  Slide 4 done")

# ── S5: Ablation ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Experiment 2 — Ablation Study: Where Do Symbolic Rules Help?",
              "Key finding: post-inference annotation adds 100% GRC coverage at zero accuracy cost")
add_img(s, charts["ablation"], 0.2, 1.35, 12.9, 5.5)
print("  Slide 5 done")

# ── S6: GRC metrics ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "GRC Metrics — NFCRM-1:2025 Compliance Outcomes",
              "100 test samples | 85 RiskCases | SHACL violations = 0")
add_img(s, charts["grc"], 0.2, 1.35, 12.9, 5.5)

# stat row at bottom
for txt, col, x in [
    ("132 nodes\n423 edges", MID_BLUE, 0.3),
    ("2,944 samples/sec\n0.34 ms/sample", GREEN, 3.3),
    ("Avg reasoning\n4.0 hops deep", MID_BLUE, 6.3),
    ("0 SHACL\nviolations", GREEN, 9.3),
]:
    badge(s, txt, x, 6.6, 2.7, 0.7, bg_col=col, size=11)
print("  Slide 6 done")

# ── S7: SHACL hallucination ───────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Experiment 3 — SHACL Hallucination Reduction (LLM+SHACL)",
              "JSON Schema provides zero improvement | SHACL reduces hallucinations to 0% among validated outputs")
add_img(s, charts["shacl"], 0.2, 1.35, 12.9, 5.5)
print("  Slide 7 done")

# ── S8: SLM CIC-IDS2018 ───────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Experiment 4 — SLM Natural-Language Classification (CIC-IDS2018)",
              "NL conversion layer: raw numbers -> plain English -> LLM classifies | 0% -> 55% accuracy")
add_img(s, charts["slm_cicids"], 0.2, 1.35, 12.9, 5.5)
print("  Slide 8 done")

# ── S9: SLM NSL-KDD ───────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Experiment 5 — SLM on NSL-KDD (T2.2 Cross-Dataset Validation)",
              "DoS 100% | Probe 100% | Overall 74% on 50-sample stratified subset")
add_img(s, charts["slm_nslkdd"], 0.2, 1.35, 12.9, 5.5)
print("  Slide 9 done")

# ── S10: ARG ──────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Asset Relationship Graph (ARG) — Multi-Source Knowledge Graph",
              "132 nodes | 423 edges | 32 real CVEs from NVD API | 5 integrated data sources")
add_img(s, charts["arg"], 0.2, 1.35, 12.9, 5.5)
print("  Slide 10 done")

# ── S11: Capability matrix ─────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Capability Matrix — What Each Approach Can and Cannot Produce",
              "The intersymbolic pipeline is the only approach producing all 6 GRC capabilities")
add_img(s, charts["capability"], 0.2, 1.35, 12.9, 5.5)
print("  Slide 11 done")

# ── S12: Effort timeline ──────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Project Effort — 573 Commits | 52 Active Days | Feb 9 – May 3, 2026",
              "9 development phases | 195+ PRs | 12 weeks of continuous iteration")
add_img(s, charts["effort"], 0.2, 1.35, 12.9, 5.5)
print("  Slide 12 done")

# ── S13: Summary ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
bg(s); header(s, "Summary — What Was Built and What It Shows")

rows = [
    (GREEN,    "RQ1 AFFIRMED",   "SHACL ontology + ARG — 100% NFCRM-1:2025 coverage | 57.6% ISO 27005 | 132 nodes, 423 edges"),
    (AMBER,    "RQ2 PARTIAL",    "Post-inference annotation: 100% GRC at zero accuracy cost | Pre-inference +0.28% (NS) | In-inference: future work"),
    (GREEN,    "RQ3 AFFIRMED",   "Capability matrix proves intersymbolic = unique GRC outputs | SHACL: 39.1% -> 0% hallucination"),
    (MID_BLUE, "SLM RESULT",     "NL conversion layer: 0% -> 55% (CIC-IDS2018) | 74% (NSL-KDD) | DoS/Probe 100% on NSL-KDD"),
    (RED,      "REMAINING",      "Human expert evaluation (3 GRC practitioners) — protocol ready in human_eval/"),
]
y = 1.4
for col, label_txt, detail in rows:
    rb = s.shapes.add_shape(1, Inches(0.3), Inches(y), Inches(12.7), Inches(0.92))
    rb.fill.solid(); rb.fill.fore_color.rgb = WHITE; rb.line.fill.background()
    badge(s, label_txt, 0.35, y+0.22, 1.7, 0.42, bg_col=col, size=11)
    add_txt(s, detail, 2.2, y+0.18, 10.6, 0.55, size=12.5, color=DGREY)
    y += 1.02

add_txt(s, "Thesis ~90% complete  |  Target submission: July 2026  |  Defense: August 2026",
        0.3, 6.65, 12.7, 0.5, size=13, bold=True,
        color=DARK_BLUE, align=PP_ALIGN.CENTER)
print("  Slide 13 done")

# ── Save ───────────────────────────────────────────────────────────────────
prs.save(OUT)
print(f"\nPresentation saved: {OUT}")
print(f"Total slides: {len(prs.slides)}")
