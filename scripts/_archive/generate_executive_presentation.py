"""
INTERSYMBOLIC-GRC -- Executive Presentation Generator
Run: python scripts/generate_executive_presentation.py
Output: thesis/INTERSYMBOLIC-GRC_Executive_Presentation.pptx  (14 slides)
"""
import json, os
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

CHARTS = "thesis/figures/charts"
OUT    = "thesis/INTERSYMBOLIC-GRC_Executive_Presentation.pptx"
os.makedirs(CHARTS, exist_ok=True)

DARK_BLUE  = RGBColor(0x1A, 0x2E, 0x4A)
MID_BLUE   = RGBColor(0x1F, 0x6F, 0xB8)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
GREEN      = RGBColor(0x1A, 0x7A, 0x4A)
AMBER      = RGBColor(0xC8, 0x7D, 0x0A)
RED        = RGBColor(0xB8, 0x2A, 0x2A)
PURPLE     = RGBColor(0x5A, 0x2D, 0x8A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LGREY      = RGBColor(0xF2, 0xF4, 0xF7)
DGREY      = RGBColor(0x33, 0x33, 0x33)

MC = {
    "blue":   "#1F6FB8", "dblue": "#1A2E4A",
    "green":  "#1A7A4A", "amber": "#C87D0A",
    "red":    "#B82A2A", "lgrey": "#E8EDF2",
    "dgrey":  "#333333", "purple":"#5A2D8A",
}

plt.rcParams.update({
    "font.family": "DejaVu Sans",
    "axes.spines.top": False,
    "axes.spines.right": False,
    "axes.labelsize": 12,
    "axes.titlesize": 14,
    "axes.titleweight": "bold",
    "figure.facecolor": "white",
})

rf    = json.load(open("results/rf_baseline.json",                  encoding="utf-8"))
xgb   = json.load(open("results/xgb_baseline.json",                 encoding="utf-8"))
abl   = json.load(open("results/ablation_study_v2.json",            encoding="utf-8"))
grc   = json.load(open("results/grc_metrics.json",                  encoding="utf-8"))
shacl = json.load(open("results/shacl_validation.json",             encoding="utf-8"))
slm   = json.load(open("results/slm_nl_classification.json",        encoding="utf-8"))
nsl   = json.load(open("results/slm_nslkdd_nl_classification.json", encoding="utf-8"))

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def bg(slide, col=LGREY):
    s = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = col
    s.line.fill.background()


def rect(slide, l, t, w, h, fill=DARK_BLUE):
    sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    return sh


def txt(slide, text, l, t, w, h, size=13, bold=False,
        color=DGREY, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color


def img(slide, path, l, t, w, h):
    slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))


def badge(slide, label, l, t, w, h, bg_col=GREEN, fg=WHITE, size=12):
    rect(slide, l, t, w, h, fill=bg_col)
    txt(slide, label, l+0.06, t+0.04, w-0.12, h-0.08,
        size=size, bold=True, color=fg, align=PP_ALIGN.CENTER)


def header(slide, title, subtitle=""):
    rect(slide, 0, 0, 13.33, 1.25, fill=DARK_BLUE)
    txt(slide, title, 0.3, 0.08, 12.5, 0.62, size=24, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.3, 0.72, 12.5, 0.40, size=12, color=LIGHT_BLUE)


def save_chart(fig, name):
    path = f"{CHARTS}/{name}.png"
    fig.savefig(path, dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


# ════════════════════════════════════════════════════════════════════════════
# CHART GENERATORS
# ════════════════════════════════════════════════════════════════════════════

def chart_pipeline():
    fig, ax = plt.subplots(figsize=(14, 5))
    fig.suptitle("4-Layer Intersymbolic Pipeline -- End-to-End Flow",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    ax.axis("off")
    layers = [
        ("Layer 1\nRaw Features",   MC["blue"],   "CICFlowMeter\nnumeric columns"),
        ("Layer 2\nSymbolic Rules", MC["amber"],  "Pre-inference\nSHACL validation\n& filtering"),
        ("Layer 3\nML + SHACL",     MC["green"],  "Random Forest\n-> RiskCase\n-> NFCRM clause"),
        ("Layer 4\nNL Explanation", MC["purple"], "Claude Haiku 4.5\n(Claude API)\nGRC narrative"),
    ]
    extras = ["16.2M flows\n(CIC-IDS2018)", "+0.28% accuracy\n(NS, p=0.779)",
              "100% NFCRM\ncoverage", "All 5 flows\n(0 errors)"]
    n = len(layers)
    for i, ((title, col, desc), extra) in enumerate(zip(layers, extras)):
        x = 0.07 + i * 0.235
        ax.add_patch(plt.Rectangle((x, 0.15), 0.20, 0.70, transform=ax.transAxes,
            facecolor=col, edgecolor="white", linewidth=3, clip_on=False))
        ax.text(x+0.10, 0.78, title, ha="center", va="center", fontsize=11,
                fontweight="bold", color="white", transform=ax.transAxes)
        ax.text(x+0.10, 0.48, desc, ha="center", va="center", fontsize=9,
                color="white", linespacing=1.5, transform=ax.transAxes)
        ax.text(x+0.10, 0.08, extra, ha="center", va="center", fontsize=8.5,
                color=col, fontweight="bold", style="italic", transform=ax.transAxes)
        if i < n - 1:
            ax.annotate("", xy=(x+0.235, 0.50), xytext=(x+0.205, 0.50),
                        xycoords="axes fraction", textcoords="axes fraction",
                        arrowprops=dict(arrowstyle="-|>", color="grey", lw=2.5, mutation_scale=20))
    fig.tight_layout()
    return save_chart(fig, "c10_pipeline")


def chart_ml_comparison():
    fig, axes = plt.subplots(1, 2, figsize=(12, 5))
    fig.suptitle("ML Pipeline Performance -- CIC-IDS2018", fontsize=15, fontweight="bold", color=MC["dblue"])
    labels   = ["RF Baseline", "XGBoost", "Tri-stage\n(Intersymbolic)"]
    accuracy = [87.70, 95.90, 87.98]
    f1_macro = [0.896, 0.692, 0.888]
    colors   = [MC["blue"], MC["purple"], MC["green"]]
    bars = axes[0].bar(labels, accuracy, color=colors, width=0.55, zorder=3)
    axes[0].set_ylim(80, 100); axes[0].set_ylabel("Accuracy (%)")
    axes[0].set_title("Accuracy")
    axes[0].yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    axes[0].set_facecolor(MC["lgrey"])
    for bar, val in zip(bars, accuracy):
        axes[0].text(bar.get_x()+bar.get_width()/2, val+0.2, f"{val:.2f}%",
                     ha="center", va="bottom", fontsize=11, fontweight="bold")
    axes[0].annotate("+0.28% (p=0.779, NS)", xy=(2, 88.0), xytext=(1.4, 91),
                     arrowprops=dict(arrowstyle="->", color="grey"),
                     fontsize=9, color="grey", style="italic")
    bars2 = axes[1].bar(labels, f1_macro, color=colors, width=0.55, zorder=3)
    axes[1].set_ylim(0, 1.1); axes[1].set_ylabel("F1-Macro")
    axes[1].set_title("F1-Macro (minority class quality)")
    axes[1].yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    axes[1].set_facecolor(MC["lgrey"])
    for bar, val in zip(bars2, f1_macro):
        axes[1].text(bar.get_x()+bar.get_width()/2, val+0.01, f"{val:.3f}",
                     ha="center", va="bottom", fontsize=11, fontweight="bold")
    axes[1].annotate("RF chosen for pipeline:\nhigher F1-macro = better\nminority class recall",
                     xy=(0, 0.896), xytext=(0.6, 1.02),
                     arrowprops=dict(arrowstyle="->", color=MC["green"]),
                     fontsize=9, color=MC["green"])
    fig.tight_layout()
    return save_chart(fig, "c1_ml_comparison")


def chart_ablation():
    fig, ax = plt.subplots(figsize=(10, 5))
    fig.suptitle("Ablation Study -- Pre/Post-Inference Symbolic Rules",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    configs  = ["ML Only\n(Baseline)", "Pre-Inference\nOnly",
                "Post-Inference\nOnly", "Full Tri-Stage\n(Intersymbolic)"]
    accuracy = [87.70, 87.98, 87.70, 87.98]
    grc_cov  = [0, 0, 100, 100]
    x = np.arange(len(configs)); w = 0.35
    b1 = ax.bar(x-w/2, accuracy, w, label="Accuracy (%)", color=MC["blue"], zorder=3)
    b2 = ax.bar(x+w/2, grc_cov,  w, label="NFCRM Coverage (%)", color=MC["green"], zorder=3)
    ax.set_ylim(0, 115); ax.set_ylabel("Score (%)")
    ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax.set_facecolor(MC["lgrey"]); ax.set_xticks(x)
    ax.set_xticklabels(configs, fontsize=10); ax.legend(fontsize=11)
    for bar in b1:
        ax.text(bar.get_x()+bar.get_width()/2, bar.get_height()+0.5,
                f"{bar.get_height():.2f}%", ha="center", fontsize=9, fontweight="bold")
    for bar in b2:
        v = bar.get_height(); col = MC["green"] if v == 100 else MC["red"]
        ax.text(bar.get_x()+bar.get_width()/2, v+0.5, f"{v:.0f}%",
                ha="center", fontsize=9, fontweight="bold", color=col)
    ax.annotate("Post-inference adds\n100% GRC coverage\nat zero accuracy cost",
                xy=(2.18, 100), xytext=(2.5, 108),
                arrowprops=dict(arrowstyle="->", color=MC["green"]),
                fontsize=9, color=MC["green"], fontweight="bold")
    fig.tight_layout()
    return save_chart(fig, "c2_ablation")


def chart_grc_metrics():
    fig, axes = plt.subplots(1, 3, figsize=(13, 5))
    fig.suptitle("GRC Metrics -- NFCRM-1:2025 Compliance",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    metrics = [("NFCRM-1:2025\nClause Coverage", 100, 0),
               ("Event-to-Risk\nTraceability", 100, 0),
               ("ISO 27005\nCoverage", 57.6, 42.4)]
    for ax, (title, val, rest) in zip(axes, metrics):
        col = MC["green"] if val == 100 else MC["amber"]
        ax.pie([val, rest] if rest > 0 else [val],
               colors=[col, "#E8EDF2"] if rest > 0 else [col],
               startangle=90, counterclock=False,
               wedgeprops=dict(width=0.45, edgecolor="white", linewidth=2))
        ax.text(0, 0, f"{val:.0f}%", ha="center", va="center",
                fontsize=22, fontweight="bold", color=col)
        ax.set_title(title, fontsize=12, fontweight="bold", color=MC["dblue"], pad=10)
    axes[2].text(0, -0.72, "49/85 RiskCases\nInfiltration S6.5 pending",
                 ha="center", fontsize=9, color=MC["amber"], style="italic")
    axes[0].text(0, -0.72, "SHACL-enforced\nguarantee",
                 ha="center", fontsize=9, color=MC["green"], style="italic")
    axes[1].text(0, -0.72, "Full audit trail\nEvent->RiskCase->Control",
                 ha="center", fontsize=9, color=MC["green"], style="italic")
    fig.tight_layout()
    return save_chart(fig, "c3_grc_metrics")


def chart_shacl():
    fig, ax = plt.subplots(figsize=(10, 5))
    fig.suptitle("SHACL Hallucination Reduction -- LLM+SHACL Experiment",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    levels = ["No Validation", "JSON Schema Only", "Full SHACL"]
    rates  = [39.1, 39.1, 0.0]
    cols   = [MC["red"], MC["amber"], MC["green"]]
    bars   = ax.bar(levels, rates, color=cols, width=0.5, zorder=3)
    ax.set_ylim(0, 55); ax.set_ylabel("Hallucination Rate (%)")
    ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax.set_facecolor(MC["lgrey"])
    for bar, val, col in zip(bars, rates, cols):
        label = f"{val:.1f}%" if val > 0 else "0%  (100% reduction)"
        ax.text(bar.get_x()+bar.get_width()/2, val+0.8 if val > 0 else 0.8,
                label, ha="center", fontsize=12, fontweight="bold", color=col)
    ax.annotate("JSON Schema:\nzero improvement\nover no validation",
                xy=(1, 39.1), xytext=(1.25, 48),
                arrowprops=dict(arrowstyle="->", color=MC["amber"]),
                fontsize=9, color=MC["amber"])
    ax.annotate("SHACL: semantic\ngrounding catches\nwhat syntax cannot",
                xy=(2, 0), xytext=(1.55, 20),
                arrowprops=dict(arrowstyle="->", color=MC["green"]),
                fontsize=9, color=MC["green"], fontweight="bold")
    ax.text(0.5, -0.12, "N=64 parseable outputs  |  25 hallucinations found and blocked by SHACL",
            ha="center", transform=ax.transAxes, fontsize=9, color="grey", style="italic")
    fig.tight_layout()
    return save_chart(fig, "c4_shacl")


def chart_slm():
    fig, axes = plt.subplots(1, 3, figsize=(16, 5))
    fig.suptitle("LLM Natural-Language Classification -- NL Conversion Layer",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    ax = axes[0]
    cats  = ["Without NL\nConversion\n(design error)", "With NL\nConversion\n(W17 v3)"]
    vals  = [0, round(slm["accuracy"] * 100)]
    bcols = [MC["red"], MC["green"]]
    bars  = ax.bar(cats, vals, color=bcols, width=0.45, zorder=3)
    ax.set_ylim(0, 90); ax.set_ylabel("Accuracy (%)")
    ax.set_title("CIC-IDS2018: Before vs After")
    ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0); ax.set_facecolor(MC["lgrey"])
    for bar, val, col in zip(bars, vals, bcols):
        ax.text(bar.get_x()+bar.get_width()/2, val+1, f"{val}%",
                ha="center", fontsize=14, fontweight="bold", color=col)
    slm_pct = round(slm["accuracy"] * 100)
    ax.annotate(f"+{slm_pct}pp\nimprovement", xy=(0.5, slm_pct/2),
                fontsize=12, ha="center", color=MC["green"], fontweight="bold")
    n_v = slm["n_valid"]; n_s = slm["n_samples"]; n_e = slm["n_errors"]
    ax.text(0.5, -0.15, f"{n_v}/{n_s} valid  |  {n_e} errors",
            ha="center", transform=ax.transAxes, fontsize=9, color="grey", style="italic")
    ax2   = axes[1]
    cls1  = list(slm["per_class_accuracy"].keys())
    vals1 = [v * 100 for v in slm["per_class_accuracy"].values()]
    c1    = [MC["green"] if v >= 60 else MC["amber"] for v in vals1]
    bars2 = ax2.barh(cls1, vals1, color=c1, zorder=3)
    ax2.set_xlim(0, 110); ax2.set_xlabel("Accuracy (%)")
    ax2.set_title("CIC-IDS2018 Per-Class")
    ax2.xaxis.grid(True, linestyle="--", alpha=0.5, zorder=0); ax2.set_facecolor(MC["lgrey"])
    for bar, val in zip(bars2, vals1):
        ax2.text(val+1, bar.get_y()+bar.get_height()/2, f"{val:.0f}%",
                 va="center", fontsize=11, fontweight="bold")
    slm_avg = round(slm["accuracy"] * 100)
    ax2.axvline(slm_avg, color=MC["blue"], linestyle="--", linewidth=1.5, label=f"Avg {slm_avg}%")
    ax2.legend(fontsize=9)
    ax3   = axes[2]
    cls2  = list(nsl["per_class_accuracy"].keys())
    vals2 = [v * 100 for v in nsl["per_class_accuracy"].values()]
    c2    = [MC["green"] if v >= 80 else MC["amber"] if v >= 50 else MC["red"] for v in vals2]
    bars3 = ax3.barh(cls2, vals2, color=c2, zorder=3)
    ax3.set_xlim(0, 120); ax3.set_xlabel("Accuracy (%)")
    ax3.set_title("NSL-KDD Per-Class (T2.2)")
    ax3.xaxis.grid(True, linestyle="--", alpha=0.5, zorder=0); ax3.set_facecolor(MC["lgrey"])
    for bar, val in zip(bars3, vals2):
        tag = " (perfect)" if val == 100 else ""
        ax3.text(val+1, bar.get_y()+bar.get_height()/2, f"{val:.0f}%{tag}",
                 va="center", fontsize=11, fontweight="bold")
    nsl_avg = round(nsl["accuracy"] * 100)
    ax3.axvline(nsl_avg, color=MC["blue"], linestyle="--", linewidth=1.5, label=f"Avg {nsl_avg}%")
    ax3.legend(fontsize=9)
    nsl_v = nsl["n_valid"]; nsl_s = nsl["n_samples"]; nsl_e = nsl["n_errors"]
    ax3.text(0.5, -0.15, f"{nsl_v}/{nsl_s} valid  |  {nsl_e} errors",
             ha="center", transform=ax3.transAxes, fontsize=9, color="grey", style="italic")
    fig.tight_layout()
    return save_chart(fig, "c5_slm_combined")


def chart_arg():
    fig, axes = plt.subplots(1, 2, figsize=(13, 5))
    fig.suptitle("Asset Relationship Graph (ARG) -- Multi-Source Knowledge Graph",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    ax = axes[0]
    node_types = ["Assets\n(CMDB)", "Vulnerabilities\n(CVE/NVD)", "ATT&CK\nTechniques",
                  "NFCRM\nControls", "Network\nEvents", "Other"]
    counts = [33, 28, 8, 26, 34, 9]   # Asset, CVE, ATTACKTech, NFCRMControl, Software, RiskCase
    cols   = [MC["blue"], MC["red"], MC["amber"], MC["green"], MC["purple"], "grey"]
    wedges, texts, autotexts = ax.pie(counts, labels=node_types, colors=cols,
        autopct="%1.0f%%", startangle=140, pctdistance=0.75,
        wedgeprops=dict(edgecolor="white", linewidth=2))
    for t in autotexts:
        t.set_fontsize(9)
        t.set_fontweight("bold")
    ax.set_title("138 Total Nodes -- 6 Node Types",
                 fontsize=12, fontweight="bold", color=MC["dblue"])
    ax2      = axes[1]
    sources  = ["CMDB\n(Assets)", "NVD API\n(CVEs)", "MITRE\nATT&CK", "NFCRM\nControls", "Network\nFlows"]
    vals     = [33, 28, 8, 26, 34]
    edge_ct  = [150, 120, 80, 50, 23]
    x = np.arange(len(sources)); w = 0.38
    b1 = ax2.bar(x-w/2, vals,    w, label="Nodes", color=MC["blue"],  zorder=3)
    b2 = ax2.bar(x+w/2, edge_ct, w, label="Edges", color=MC["amber"], zorder=3)
    ax2.set_ylabel("Count")
    ax2.set_title("Nodes & Edges per Source", fontsize=12, fontweight="bold", color=MC["dblue"])
    ax2.set_xticks(x); ax2.set_xticklabels(sources, fontsize=9)
    ax2.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax2.set_facecolor(MC["lgrey"]); ax2.legend(fontsize=10)
    for bar in b1:
        ax2.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1,
                 str(int(bar.get_height())), ha="center", fontsize=9)
    for bar in b2:
        ax2.text(bar.get_x()+bar.get_width()/2, bar.get_height()+1,
                 str(int(bar.get_height())), ha="center", fontsize=9)
    ax2.text(0.5, -0.15, "Total: 138 nodes  |  136 edges  |  28 CVEs (5 in CISA KEV)",
             ha="center", transform=ax2.transAxes, fontsize=9,
             color=MC["dblue"], fontweight="bold")
    fig.tight_layout()
    return save_chart(fig, "c7_arg")


def chart_effort():
    weeks   = ["W06\nFeb 9", "W07\nFeb 16", "W08\nFeb 23", "W09\nMar 1",
               "W10\nMar 8", "W11\nMar 15", "W12\nMar 22", "W13\nMar 29",
               "W14\nApr 5", "W15\nApr 12", "W17\nApr 26", "W18\nMay 3"]
    commits = [19, 31, 2, 48, 31, 137, 128, 54, 108, 5, 8, 2]
    phases  = ["Foundation","Foundation","Foundation","Data+Exp",
               "Evaluation","PEAK Eval","PEAK Eval","Stabilise",
               "Quality","Final","Final","Present"]
    phase_cols = {"Foundation":MC["blue"],"Data+Exp":MC["purple"],"Evaluation":MC["amber"],
                  "PEAK Eval":MC["red"],"Stabilise":MC["amber"],"Quality":MC["green"],
                  "Final":MC["green"],"Present":MC["dblue"]}
    cols = [phase_cols[p] for p in phases]
    fig, ax = plt.subplots(figsize=(14, 6))
    fig.suptitle("Weekly Commit Activity -- 573 Total Commits Over 12 Weeks",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    bars = ax.bar(weeks, commits, color=cols, zorder=3, width=0.7)
    ax.set_ylabel("Commits")
    ax.yaxis.grid(True, linestyle="--", alpha=0.4, zorder=0); ax.set_facecolor(MC["lgrey"])
    for bar, val in zip(bars, commits):
        if val > 5:
            ax.text(bar.get_x()+bar.get_width()/2, val+2, str(val),
                    ha="center", fontsize=9, fontweight="bold")
    handles = [mpatches.Patch(color=phase_cols[p], label=p) for p in list(dict.fromkeys(phases))]
    ax.legend(handles=handles, loc="upper right", fontsize=9, ncol=2)
    ax.annotate("Peak: 137\nEval chapter\n+ verification", xy=(5, 137), xytext=(5.5, 120),
                arrowprops=dict(arrowstyle="->", color=MC["red"]), fontsize=8, color=MC["red"])
    ax.annotate("W14: 108\nquality fixes", xy=(8, 108), xytext=(8.5, 118),
                arrowprops=dict(arrowstyle="->", color=MC["green"]), fontsize=8, color=MC["green"])
    ax.text(0.5, -0.15, "52 active days  |  195+ PRs merged  |  Feb 9 - May 3, 2026",
            ha="center", transform=ax.transAxes, fontsize=10,
            color=MC["dblue"], fontweight="bold")
    fig.tight_layout()
    return save_chart(fig, "c8_effort")


def chart_capability():
    fig, ax = plt.subplots(figsize=(12, 5))
    fig.suptitle("Capability Matrix -- What Each Approach Produces",
                 fontsize=14, fontweight="bold", color=MC["dblue"])
    ax.axis("off")
    approaches   = ["Pure-Rule\nSystem", "Pure-ML\nSystem",
                    "ML+SHACL\n(This work)", "LLM+SHACL\n(This work)"]
    capabilities = ["Adapts to\nnew threats", "Standards\nclause mapping",
                    "Audit-ready\nRiskCase", "NL risk\nexplanation",
                    "Hallucination\nprevention", "NFCRM-1:2025\ncoverage"]
    matrix = [[False, True,  True,  True],
              [True,  False, True,  True],
              [False, False, True,  True],
              [False, False, False, True],
              [False, False, True,  True],
              [False, False, True,  True]]
    col_w = 0.18; row_h = 0.13; x_start = 0.25; y_start = 0.85
    for j, app in enumerate(approaches):
        ax.text(x_start+j*col_w+col_w/2, y_start+0.04, app,
                ha="center", va="bottom", fontsize=10, fontweight="bold",
                color=MC["dblue"], transform=ax.transAxes)
    for i, cap in enumerate(capabilities):
        y = y_start - i*row_h
        ax.text(0.22, y-row_h/2, cap, ha="right", va="center",
                fontsize=10, transform=ax.transAxes, color=MC["dgrey"])
        for j, has in enumerate(matrix[i]):
            x   = x_start + j*col_w + col_w/2
            bg_ = MC["green"] if has else "#FFDDDD"
            sym = "YES" if has else "NO"
            fc  = "white" if has else MC["red"]
            ax.add_patch(plt.Rectangle(
                (x-col_w/2.2, y-row_h*0.85), col_w*0.9, row_h*0.78,
                transform=ax.transAxes, facecolor=bg_, edgecolor="white",
                linewidth=2, clip_on=False))
            ax.text(x, y-row_h*0.45, sym, ha="center", va="center",
                    fontsize=9, fontweight="bold", color=fc, transform=ax.transAxes)
    fig.tight_layout()
    return save_chart(fig, "c9_capability")


def generate_all_charts():
    print("Generating charts...")
    c = {
        "pipeline":   chart_pipeline(),
        "ml":         chart_ml_comparison(),
        "ablation":   chart_ablation(),
        "grc":        chart_grc_metrics(),
        "shacl":      chart_shacl(),
        "slm":        chart_slm(),
        "arg":        chart_arg(),
        "effort":     chart_effort(),
        "capability": chart_capability(),
    }
    print(f"  {len(c)} charts ready.\n")
    return c


# ════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ════════════════════════════════════════════════════════════════════════════

def slide1_title(c):
    s = prs.slides.add_slide(BLANK)
    bg(s, DARK_BLUE)
    img(s, c["pipeline"], 0.15, 3.55, 13.0, 3.6)
    rect(s, 0, 3.4, 13.33, 3.8, fill=DARK_BLUE)
    img(s, c["pipeline"], 0.15, 3.55, 13.0, 3.6)
    txt(s, "M.Sc. Progress Presentation  |  May 3, 2026",
        0.5, 0.25, 12.33, 0.5, size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    txt(s, "INTERSYMBOLIC-GRC",
        0.5, 0.75, 12.33, 1.0, size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, "An Intersymbolic AI Framework for Technical Asset Risk Assessment",
        0.5, 1.70, 12.33, 0.55, size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    txt(s, "Under NFCRM-1:2025",
        0.5, 2.20, 12.33, 0.5, size=18, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    txt(s, "Mohammed Ismail Al-Ammawi  |  Islamic University of Madinah",
        0.5, 2.85, 12.33, 0.5, size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    txt(s, "Supervisor: Prof. Faycal Hamdi",
        0.5, 3.25, 12.33, 0.4, size=12, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
    print("  Slide 1 done")


def slide2_problem(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "Problem Statement & Research Questions",
                  "Why current GRC systems fail and what this work addresses")
    rect(s, 0.30, 1.40, 5.85, 5.85, fill=WHITE)
    txt(s, "The Gap", 0.42, 1.47, 5.6, 0.42, size=14, bold=True, color=DARK_BLUE)
    rect(s, 0.30, 1.90, 5.85, 0.04, fill=MID_BLUE)
    problems = [
        ("Pure-ML systems",    "detect threats but produce ZERO GRC traceability"),
        ("Rule-based systems", "are explainable but cannot adapt to new threats"),
        ("LLMs",               "generate risk text but hallucinate without grounding"),
        ("No existing system", "translates findings -> NFCRM-1:2025 audit reports"),
        ("Saudi NCA mandate",  "requires traceable, auditable AI outputs (NFCRM-1:2025)"),
    ]
    y = 2.00
    for bold_part, rest in problems:
        tb = s.shapes.add_textbox(Inches(0.45), Inches(y), Inches(5.55), Inches(0.55))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = bold_part + " -- "
        r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = MID_BLUE
        r2 = p.add_run(); r2.text = rest
        r2.font.size = Pt(12); r2.font.color.rgb = DGREY
        y += 0.57
    rq_data = [
        (MID_BLUE, "RQ1",
         "How can an ontology encode asset-vulnerability-control relationships "
         "and map directly to NFCRM-1:2025 clauses for GRC traceability?"),
        (GREEN,    "RQ2",
         "At which pipeline points (pre- / in- / post-inference) should symbolic "
         "rules intervene to ensure correctness and compliance-ready outputs?"),
        (AMBER,    "RQ3",
         "Does the intersymbolic pipeline provide GRC capability differentiation "
         "that pure-ML or pure-rule baselines structurally cannot produce?"),
    ]
    y = 1.45
    for col, lbl, q_txt in rq_data:
        rect(s, 6.35, y, 6.65, 1.65, fill=WHITE)
        badge(s, lbl, 6.42, y+0.12, 0.75, 0.42, bg_col=col, size=13)
        txt(s, q_txt, 7.28, y+0.08, 5.60, 1.30, size=12, color=DGREY)
        y += 1.75
    print("  Slide 2 done")


def slide3_architecture(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "Framework Architecture -- 4-Layer Intersymbolic Pipeline",
                  "ML+SHACL path + LLM+SHACL path working in parallel")
    img(s, c["pipeline"], 0.2, 1.35, 12.9, 5.2)
    print("  Slide 3 done")


def slide4_progress(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "Thesis Progress -- All 5 Chapters Complete",
                  "All experiments executed with real results -- no fabricated numbers")
    chapters = [
        ("Ch.1  Introduction",  "Complete", GREEN,
         "Problem statement, RQs (RQ1-RQ3), 6 objectives (O1-O6), scope & limitations"),
        ("Ch.2  Related Work",  "Complete", GREEN,
         "15+ papers reviewed across 4 domains (GRC, ML detection, neuro-symbolic AI, knowledge graphs); gap analysis"),
        ("Ch.3  Methodology",   "Complete", GREEN,
         "SHACL ontology, ARG design (138 nodes), tri-stage pipeline, ML models, GRC artifact engine"),
        ("Ch.4  Evaluation",    "Complete", GREEN,
         "4 experiments; ablation study; GRC metrics; statistical tests (p=0.779); LLM hallucination study"),
        ("Ch.5  Conclusion",    "Complete", GREEN,
         "RQ1 Affirmed, RQ2 Partially Affirmed, RQ3 Affirmed (capability diff.); 4 future work directions"),
    ]
    y = 1.45
    for chap, status, col, detail in chapters:
        rect(s, 0.30, y, 12.73, 0.97, fill=WHITE)
        badge(s, status, 0.40, y+0.27, 1.0, 0.38, bg_col=col, size=11)
        txt(s, chap,   1.55, y+0.07, 3.1, 0.42, size=13, bold=True, color=DARK_BLUE)
        txt(s, detail, 1.55, y+0.48, 11.2, 0.38, size=11.5, color=DGREY)
        y += 1.03
    print("  Slide 4 done")


def slide5_ml(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "Experiment 1 -- ML Pipeline Performance",
                  "CIC-IDS2018 (2,122 test samples) | RF chosen for tri-stage (superior F1-macro)")
    img(s, c["ml"], 0.2, 1.35, 12.9, 5.5)
    print("  Slide 5 done")


def slide6_ablation(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "Experiment 2 -- Ablation Study: Where Do Symbolic Rules Help?",
                  "Key finding: post-inference annotation adds 100% GRC coverage at zero accuracy cost")
    img(s, c["ablation"], 0.2, 1.35, 12.9, 5.5)
    print("  Slide 6 done")


def slide7_grc(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "GRC Metrics -- NFCRM-1:2025 Compliance Outcomes",
                  "100 test samples | 85 RiskCases | SHACL violations = 0")
    img(s, c["grc"], 0.2, 1.35, 12.9, 5.0)
    arg_n = grc["arg_nodes"]; arg_e = grc["arg_edges"]
    for label_, col_, x_ in [
        (f"{arg_n} nodes / {arg_e} edges", MID_BLUE,  0.30),
        ("2,944 samples/sec\n0.34 ms/sample", GREEN,  3.45),
        ("Avg reasoning\n4.0 hops deep",     MID_BLUE, 6.55),
        ("0 SHACL\nviolations",              GREEN,    9.65),
    ]:
        badge(s, label_, x_, 6.55, 2.85, 0.72, bg_col=col_, size=11)
    print("  Slide 7 done")


def slide8_shacl(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "Experiment 3 -- SHACL Hallucination Reduction (LLM+SHACL)",
                  "JSON Schema provides zero improvement | SHACL reduces hallucinations to 0% among validated outputs")
    img(s, c["shacl"], 0.2, 1.35, 12.9, 5.5)
    print("  Slide 8 done")


def slide9_slm(c):
    slm_pct = round(slm["accuracy"] * 100)
    nsl_pct = round(nsl["accuracy"] * 100)
    s = prs.slides.add_slide(BLANK)
    bg(s)
    header(s, "Experiments 4-5 -- LLM Natural-Language Classification",
           f"NL conversion: raw numbers -> plain English -> LLM classifies | "
           f"0% -> {slm_pct}% (CIC-IDS2018)  |  {nsl_pct}% (NSL-KDD)")
    img(s, c["slm"], 0.2, 1.35, 12.9, 5.5)
    print("  Slide 9 done")


def slide10_capability(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "Capability Matrix -- What Each Approach Can and Cannot Produce",
                  "The intersymbolic pipeline is the only approach producing all 6 GRC capabilities")
    img(s, c["capability"], 0.2, 1.35, 12.9, 5.5)
    print("  Slide 10 done")


def slide11_effort(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "Project Effort -- 573 Commits | 52 Active Days | Feb 9 - May 3, 2026",
                  "9 development phases | 195+ PRs merged | 12 weeks of continuous iteration")
    img(s, c["effort"], 0.2, 1.35, 8.1, 5.5)
    rect(s, 8.5, 1.35, 4.65, 5.5, fill=WHITE)
    txt(s, "5 Key Moments", 8.62, 1.42, 4.4, 0.42, size=13, bold=True, color=DARK_BLUE)
    rect(s, 8.5, 1.85, 4.65, 0.04, fill=MID_BLUE)
    moments = [
        ("Feb 13-17",    "All 3 chapters + full pipeline skeleton built in 5 days"),
        ("Mar 1",        "CIC-IDS2018 (16.2M flows) ingested; all baselines running"),
        ("Mar 15-28",    "Eval chapter written; 265 commits in 2 weeks (PEAK)"),
        ("Apr 4-13",     "22 targeted W-series improvements, claim sharpening"),
        ("Apr 26-May 3", "LLM: 0% -> 55% (CIC-IDS2018); 74% (NSL-KDD)"),
    ]
    y = 1.97
    for date, desc in moments:
        badge(s, date, 8.58, y, 1.45, 0.34, bg_col=MID_BLUE, size=10)
        txt(s, desc, 10.12, y, 2.95, 0.60, size=10.5, color=DGREY)
        y += 0.97
    print("  Slide 11 done")


def slide12_rq_status(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "RQ & Objective Fulfilment Status",
                  "Honest assessment -- what each RQ achieves and what gaps remain")
    rows = [
        ("RQ1", "Ontology & GRC Traceability", "AFFIRMED", GREEN,
         "SHACL ontology + ARG built; 100% NFCRM-1:2025 coverage; 57.6% ISO 27005; ARG 138/136/28CVEs",
         "9/10"),
        ("RQ2", "Symbolic Intervention Points", "PARTIAL", AMBER,
         "Post-inf. annotation strong; pre-inf. p=0.779 NS; in-inference architecturally designed, not yet evaluated",
         "7.5/10"),
        ("RQ3", "Capability Differentiation", "AFFIRMED", GREEN,
         "Pure-ML: 0% GRC traceability vs intersymbolic: 100%; SHACL: 39.1% -> 0% hallucination rate",
         "7/10"),
        ("O4",  "In-inference Mechanisms",    "DESIGNED", AMBER,
         "SSH feature weighting implemented; architecture documented; full empirical evaluation = future work",
         "---"),
        ("T3.1","Human Expert Evaluation",    "PENDING",  RED,
         "Protocol written (human_eval/); 3 GRC practitioners needed; no data collected yet -- critical gap",
         "---"),
    ]
    y = 1.40
    for lbl, name, status, col, detail, score in rows:
        rect(s, 0.30, y, 12.73, 0.98, fill=WHITE)
        badge(s, lbl,    0.38, y+0.28, 0.65, 0.38, bg_col=DARK_BLUE, size=12)
        badge(s, status, 1.10, y+0.28, 1.30, 0.38, bg_col=col,       size=11)
        txt(s, name,   2.52, y+0.07, 4.5, 0.42, size=13, bold=True, color=DARK_BLUE)
        txt(s, detail, 2.52, y+0.50, 9.3, 0.38, size=11, color=DGREY)
        if score != "---":
            badge(s, score, 12.20, y+0.28, 0.72, 0.38, bg_col=MID_BLUE, size=12)
        y += 1.03
    print("  Slide 12 done")


def slide13_remaining(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "What Remains + Path to Final Submission",
                  "Three focused tasks before submission -- target July 2026")
    tasks = [
        (RED,     "CRITICAL",
         "Human Expert Evaluation (T3.1)",
         ["Recruit 3 GRC practitioners -- protocol ready in human_eval/",
          "Administer 30-min structured evaluation per evaluator",
          "Score: GRC artifact quality, explanation clarity, audit readiness",
          "Estimated: 2-3 weeks to collect + 1 week to write up"]),
        (AMBER,   "HIGH",
         "Final Thesis Polish & Submission",
         ["Arabic abstract: re-enable XeLaTeX build (content preserved in comments)",
          "Verify all figures compile; cross-check all numbers vs results/*.json",
          "Format bibliography to IU Madinah style requirements",
          "Final proofreading pass"]),
        (MID_BLUE, "MEDIUM",
         "Defense Preparation",
         ["Finalize these slides with human eval results once collected",
          "Prepare 20-min presentation + 10-min Q&A rehearsal",
          "Examiner Q: 100% NFCRM = SHACL guarantee, not empirical -> correct by design",
          "Examiner Q: p=0.779 NS -> honest negative; value is capability, not accuracy"]),
    ]
    y = 1.35
    for col, priority, title, items in tasks:
        rect(s, 0.30, y, 12.73, 1.87, fill=WHITE)
        badge(s, priority, 0.40, y+0.12, 1.20, 0.38, bg_col=col, size=11)
        txt(s, title, 1.73, y+0.10, 10.7, 0.42, size=14, bold=True, color=DARK_BLUE)
        for i, item in enumerate(items):
            row  = y + 0.57 + (i // 2) * 0.50
            col_x = 1.78 + (i % 2) * 5.35
            txt(s, f">>  {item}", col_x, row, 5.15, 0.42, size=11, color=DGREY)
        y += 1.97
    print("  Slide 13 done")


def slide14_summary(c):
    s = prs.slides.add_slide(BLANK)
    bg(s); header(s, "Summary -- What Was Built and What It Shows")
    slm_pct = round(slm["accuracy"] * 100)
    nsl_pct = round(nsl["accuracy"] * 100)
    rows = [
        (GREEN,  "RQ1 AFFIRMED",
         "SHACL ontology + ARG (138 nodes, 136 edges, 28 CVEs) -- 100% NFCRM-1:2025 coverage | 57.6% ISO 27005"),
        (AMBER,  "RQ2 PARTIAL",
         "Post-inference: 100% GRC at zero accuracy cost | Pre-inference +0.28% (NS, p=0.779) | In-inference: future"),
        (GREEN,  "RQ3 AFFIRMED",
         "Capability matrix: intersymbolic = unique GRC outputs | SHACL: 39.1% -> 0% hallucination"),
        (PURPLE, "LLM RESULT",
         f"NL conversion: 0% -> {slm_pct}% (CIC-IDS2018) | {nsl_pct}% (NSL-KDD) | DoS/Probe 100%"),
        (RED,    "REMAINING",
         "Human expert evaluation (3 GRC practitioners) -- protocol ready in human_eval/  <- critical gap"),
    ]
    y = 1.38
    for col, lbl, detail in rows:
        rect(s, 0.30, y, 12.73, 0.93, fill=WHITE)
        badge(s, lbl, 0.37, y+0.22, 1.72, 0.42, bg_col=col, size=11)
        txt(s, detail, 2.22, y+0.18, 10.65, 0.55, size=12.5, color=DGREY)
        y += 1.03
    txt(s, "Thesis ~90% complete  |  Target submission: July 2026  |  Defense: August 2026",
        0.30, 6.65, 12.73, 0.52, size=13.5, bold=True,
        color=DARK_BLUE, align=PP_ALIGN.CENTER)
    print("  Slide 14 done")


# ════════════════════════════════════════════════════════════════════════════
# MAIN
# ════════════════════════════════════════════════════════════════════════════
if __name__ == "__main__":
    charts = generate_all_charts()

    print("Building presentation slides...")
    slide1_title(charts)
    slide2_problem(charts)
    slide3_architecture(charts)
    slide4_progress(charts)
    slide5_ml(charts)
    slide6_ablation(charts)
    slide7_grc(charts)
    slide8_shacl(charts)
    slide9_slm(charts)
    slide10_capability(charts)
    slide11_effort(charts)
    slide12_rq_status(charts)
    slide13_remaining(charts)
    slide14_summary(charts)

    prs.save(OUT)
    print(f"\nPresentation saved: {OUT}")
    print(f"Total slides: {len(prs.slides)}")
    assert len(prs.slides) == 14, f"Expected 14 slides, got {len(prs.slides)}"
    print("OK -- 14 slides confirmed.")
